"""Voice agent with PowerPoint slide viewer and content speaker."""

import asyncio
import os
import threading
from datetime import datetime
from pathlib import Path
from typing import Optional, List, Dict

import tkinter as tk
from tkinter import ttk, scrolledtext
from PIL import Image, ImageTk
from pptx import Presentation
from pptx.util import Inches, Pt
from loguru import logger
from dotenv import load_dotenv
from pipecat.pipeline.pipeline import Pipeline
from pipecat.pipeline.runner import PipelineRunner
from pipecat.pipeline.task import PipelineParams, PipelineTask
from pipecat.services.openai.realtime.llm import OpenAIRealtimeLLMService
from pipecat.services.llm_service import FunctionCallParams
from pipecat.transports.local.audio import LocalAudioTransport, LocalAudioTransportParams
from pipecat.audio.vad.silero import SileroVADAnalyzer
from pipecat.adapters.schemas.function_schema import FunctionSchema
from pipecat.adapters.schemas.tools_schema import ToolsSchema
from pipecat.processors.aggregators.llm_context import LLMContext
from pipecat.processors.aggregators.llm_response_universal import LLMContextAggregatorPair
from pipecat.services.openai.realtime.events import (
    AudioConfiguration,
    AudioInput,
    InputAudioTranscription,
    SemanticTurnDetection,
    SessionProperties,
)

load_dotenv()


class SlideData:
    """Container for slide information."""
    
    def __init__(self, slide_number: int, title: str, content: str, raw_text: str):
        self.slide_number = slide_number
        self.title = title
        self.content = content
        self.raw_text = raw_text


class PowerPointViewerUI:
    """Tkinter UI for displaying PowerPoint slides with text content."""
    
    def __init__(self, pptx_file: str = "./slides.pptx"):
        self.pptx_file = Path(pptx_file)
        self.current_index = 0
        self.slides: List[SlideData] = []
        self.root = None
        self.title_label = None
        self.content_text = None
        self.info_label = None
        
        # Load slides
        self._load_slides()
        
    def _extract_text_from_shape(self, shape) -> str:
        """Extract text from a shape."""
        if hasattr(shape, "text"):
            return shape.text.strip()
        elif hasattr(shape, "text_frame"):
            return "\n".join([
                paragraph.text.strip() 
                for paragraph in shape.text_frame.paragraphs
                if paragraph.text.strip()
            ])
        return ""
    
    def _load_slides(self):
        """Load all slides from the PowerPoint file."""
        if not self.pptx_file.exists():
            logger.error(f"PowerPoint file not found: {self.pptx_file}")
            return
        
        try:
            prs = Presentation(str(self.pptx_file))
            logger.info(f"Loading presentation with {len(prs.slides)} slides")
            
            for idx, slide in enumerate(prs.slides, 1):
                # Extract title
                title = "Untitled Slide"
                if slide.shapes.title:
                    title = slide.shapes.title.text.strip() or f"Slide {idx}"
                
                # Extract all text content
                content_parts = []
                raw_text_parts = []
                
                for shape in slide.shapes:
                    text = self._extract_text_from_shape(shape)
                    if text:
                        raw_text_parts.append(text)
                        # Skip title in content (avoid duplication)
                        if shape != slide.shapes.title:
                            content_parts.append(text)
                
                content = "\n\n".join(content_parts)
                raw_text = "\n\n".join(raw_text_parts)
                
                slide_data = SlideData(
                    slide_number=idx,
                    title=title,
                    content=content if content else "No content",
                    raw_text=raw_text if raw_text else "No text found"
                )
                
                self.slides.append(slide_data)
                logger.info(f"Loaded slide {idx}: {title}")
            
            logger.info(f"Successfully loaded {len(self.slides)} slides")
            
        except Exception as e:
            logger.error(f"Error loading PowerPoint: {e}")
    
    def start(self):
        """Start the Tkinter UI in a separate thread."""
        thread = threading.Thread(target=self._run_ui, daemon=True)
        thread.start()
        
    def _run_ui(self):
        """Run the Tkinter main loop."""
        self.root = tk.Tk()
        self.root.title("Voice Controlled PowerPoint Viewer")
        self.root.geometry("1000x700")
        self.root.configure(bg="#2b2b2b")
        
        # App title
        app_title = tk.Label(
            self.root,
            text="🎤 Voice Controlled PowerPoint Viewer",
            font=("Arial", 16, "bold"),
            bg="#2b2b2b",
            fg="white"
        )
        app_title.pack(pady=10)
        
        # Slide title frame
        title_frame = tk.Frame(self.root, bg="#1e1e1e", relief=tk.RAISED, borderwidth=2)
        title_frame.pack(padx=20, pady=(5, 10), fill=tk.X)
        
        self.title_label = tk.Label(
            title_frame,
            text="No slides loaded",
            font=("Arial", 20, "bold"),
            bg="#1e1e1e",
            fg="#4ECDC4",
            wraplength=900,
            justify=tk.LEFT,
            padx=20,
            pady=15
        )
        self.title_label.pack(fill=tk.X)
        
        # Content frame
        content_frame = tk.Frame(self.root, bg="#1e1e1e", relief=tk.SUNKEN, borderwidth=2)
        content_frame.pack(padx=20, pady=5, fill=tk.BOTH, expand=True)
        
        # Scrolled text widget for content
        self.content_text = scrolledtext.ScrolledText(
            content_frame,
            wrap=tk.WORD,
            font=("Arial", 12),
            bg="#1e1e1e",
            fg="white",
            padx=20,
            pady=20,
            relief=tk.FLAT,
            state=tk.DISABLED
        )
        self.content_text.pack(fill=tk.BOTH, expand=True)
        
        # Info label
        self.info_label = tk.Label(
            self.root,
            text="No slides loaded",
            font=("Arial", 10),
            bg="#2b2b2b",
            fg="#888888"
        )
        self.info_label.pack(pady=5)
        
        # Navigation buttons
        button_frame = tk.Frame(self.root, bg="#2b2b2b")
        button_frame.pack(pady=10)
        
        prev_btn = ttk.Button(button_frame, text="← Previous", command=self.previous_slide)
        prev_btn.pack(side=tk.LEFT, padx=5)
        
        next_btn = ttk.Button(button_frame, text="Next →", command=self.next_slide)
        next_btn.pack(side=tk.LEFT, padx=5)
        
        # Display first slide if available
        if self.slides:
            self.display_current_slide()
        
        self.root.mainloop()
    
    def display_current_slide(self):
        """Display the current slide."""
        if not self.slides:
            return
            
        try:
            slide = self.slides[self.current_index]
            
            # Update title
            if self.root:
                self.title_label.config(text=slide.title)
                
                # Update content
                self.content_text.config(state=tk.NORMAL)
                self.content_text.delete(1.0, tk.END)
                self.content_text.insert(1.0, slide.content)
                self.content_text.config(state=tk.DISABLED)
                
                # Update info
                self.info_label.config(
                    text=f"Slide {slide.slide_number} of {len(self.slides)}"
                )
                
            logger.info(f"Displaying slide {slide.slide_number}: {slide.title}")
            
        except Exception as e:
            logger.error(f"Error displaying slide: {e}")
    
    def next_slide(self):
        """Move to the next slide."""
        if not self.slides:
            return None
            
        self.current_index = (self.current_index + 1) % len(self.slides)
        self.display_current_slide()
        return self.get_current_slide_info()
    
    def previous_slide(self):
        """Move to the previous slide."""
        if not self.slides:
            return None
            
        self.current_index = (self.current_index - 1) % len(self.slides)
        self.display_current_slide()
        return self.get_current_slide_info()
    
    def goto_slide(self, slide_number: int):
        """Go to a specific slide by number (1-based)."""
        if not self.slides:
            return {"error": "No slides available"}
            
        # Convert to 0-based index
        index = slide_number - 1
        
        if 0 <= index < len(self.slides):
            self.current_index = index
            self.display_current_slide()
            return self.get_current_slide_info()
        else:
            return {"error": f"Invalid slide number. Please choose between 1 and {len(self.slides)}"}
    
    def get_current_slide_info(self):
        """Get information about the current slide including full text."""
        if not self.slides:
            return {"error": "No slides available"}
            
        slide = self.slides[self.current_index]
        return {
            "slide_number": slide.slide_number,
            "total_slides": len(self.slides),
            "title": slide.title,
            "content": slide.raw_text,  # Full text for AI to read
            "has_next": self.current_index < len(self.slides) - 1,
            "has_previous": self.current_index > 0
        }


# Global viewer instance
slide_viewer: Optional[PowerPointViewerUI] = None


# Function handlers
async def next_slide(params: FunctionCallParams) -> None:
    """Navigate to the next slide and return its content."""
    global slide_viewer
    
    if slide_viewer:
        info = slide_viewer.next_slide()
        logger.info(f"📊 Next slide: {info.get('title', 'N/A')}")
        await params.result_callback(info)
    else:
        await params.result_callback({"error": "Slide viewer not initialized"})


async def previous_slide(params: FunctionCallParams) -> None:
    """Navigate to the previous slide and return its content."""
    global slide_viewer
    
    if slide_viewer:
        info = slide_viewer.previous_slide()
        logger.info(f"📊 Previous slide: {info.get('title', 'N/A')}")
        await params.result_callback(info)
    else:
        await params.result_callback({"error": "Slide viewer not initialized"})


async def goto_slide(params: FunctionCallParams) -> None:
    """Go to a specific slide by number and return its content."""
    global slide_viewer
    
    slide_number = params.arguments.get("slide_number", 1)
    
    if slide_viewer:
        info = slide_viewer.goto_slide(slide_number)
        logger.info(f"📊 Go to slide {slide_number}: {info.get('title', 'N/A')}")
        await params.result_callback(info)
    else:
        await params.result_callback({"error": "Slide viewer not initialized"})


async def get_current_slide_content(params: FunctionCallParams) -> None:
    """Get the full content of the current slide for speaking."""
    global slide_viewer
    
    if slide_viewer:
        info = slide_viewer.get_current_slide_info()
        logger.info(f"📊 Current slide content requested: {info.get('title', 'N/A')}")
        await params.result_callback(info)
    else:
        await params.result_callback({"error": "Slide viewer not initialized"})


async def main() -> None:
    """Main entry point for the voice agent."""
    global slide_viewer
    
    # Check for OpenAI API key
    if not os.getenv("OPENAI_API_KEY"):
        logger.error("Missing OPENAI_API_KEY in .env file")
        logger.info("Get your API key from: https://platform.openai.com/api-keys")
        return

    # Initialize and start slide viewer UI
    logger.info("📊 Initializing PowerPoint viewer UI...")
    slide_viewer = PowerPointViewerUI("./slides.pptx")
    
    if not slide_viewer.slides:
        logger.error("No slides loaded! Please ensure slides.pptx exists.")
        return
    
    slide_viewer.start()
    
    # Give UI time to initialize
    await asyncio.sleep(1)

    logger.info("🎤 Initializing local microphone transport...")
    
    transport = LocalAudioTransport(
        params=LocalAudioTransportParams(
            audio_in_enabled=True,
            audio_out_enabled=True,
            vad_analyzer=SileroVADAnalyzer(),
        )
    )

    logger.info("🤖 Initializing OpenAI Realtime Service...")
    
    # Define function schemas
    next_slide_function = FunctionSchema(
        name="next_slide",
        description="Navigate to the next slide and get its content",
        properties={},
        required=[],
    )
    
    previous_slide_function = FunctionSchema(
        name="previous_slide",
        description="Navigate to the previous slide and get its content",
        properties={},
        required=[],
    )
    
    goto_slide_function = FunctionSchema(
        name="goto_slide",
        description="Go to a specific slide by its number (1-based index)",
        properties={
            "slide_number": {
                "type": "integer",
                "description": "The slide number to navigate to (1 for first slide, 2 for second, etc.)",
            }
        },
        required=["slide_number"],
    )
    
    get_current_slide_content_function = FunctionSchema(
        name="get_current_slide_content",
        description="Get the complete content and information about the currently displayed slide",
        properties={},
        required=[],
    )
    
    # Create tools schema
    tools = ToolsSchema(
        standard_tools=[
            next_slide_function,
            previous_slide_function,
            goto_slide_function,
            get_current_slide_content_function,
        ]
    )
    
    # Set up session properties with instructions
    session_properties = SessionProperties(
        audio=AudioConfiguration(
            input=AudioInput(
                transcription=InputAudioTranscription(),
                turn_detection=SemanticTurnDetection(),
            )
        ),
        instructions="""You are a helpful voice assistant that presents PowerPoint slides to users.

Your workflow for presenting slides:
1. When starting or moving to a new slide, ALWAYS use get_current_slide_content() to retrieve the slide information
2. Read out the slide title clearly
3. Then present the main content/points from the slide in a natural, conversational way
4. After explaining the slide content, ALWAYS ask: "Do you have any questions about this slide?" or "Would you like me to clarify anything?"
5. Wait for the user's response
6. If they say "next" or want to move on, use next_slide() and repeat the process
7. If they have questions, answer them based on the slide content
8. If they say "previous" or "go back", use previous_slide()

IMPORTANT BEHAVIOR:
- After presenting each slide, you MUST pause and explicitly ask if they have questions
- Be conversational and engaging when presenting the content
- Don't just read the text verbatim - explain it naturally
- If asked to go to a specific slide number, use goto_slide(number)
- Listen carefully for questions and provide helpful answers based on the slide content

Commands you should recognize:
- "next slide" or "next" → use next_slide()
- "previous slide" or "previous" or "go back" → use previous_slide()
- "go to slide 3" → use goto_slide(3)
- When user asks questions, answer based on the slide content without moving slides

Start by getting the first slide content and presenting it, then ask if they have questions."""
    )
    
    # Get initial slide info for the starting message
    initial_slide = slide_viewer.get_current_slide_info()
    starting_message = f"""Welcome! I'll be presenting the slides for you today. 
    
Let me start with the first slide. This slide is titled "{initial_slide['title']}".

Here's the content:
{initial_slide['content']}

Do you have any questions about this slide, or shall we move to the next one?"""

    # Initialize OpenAI Realtime Service
    realtime = OpenAIRealtimeLLMService(
        api_key=os.getenv("OPENAI_API_KEY"),
        model="gpt-4o-realtime-preview-2024-12-17",
        voice="alloy",
        session_properties=session_properties,
        start_audio_paused=False,
    )

    # Register the functions
    realtime.register_function("next_slide", next_slide)
    realtime.register_function("previous_slide", previous_slide)
    realtime.register_function("goto_slide", goto_slide)
    realtime.register_function("get_current_slide_content", get_current_slide_content)
    
    # Create LLM context with tools and initial context about the first slide
    context = LLMContext(
        messages=[
            {
                "role": "user", 
                "content": f"Please start the presentation. The first slide is titled '{initial_slide['title']}'. Here's the content: {initial_slide['content']}"
            }
        ],
        tools=tools,
    )
    
    # Create context aggregator
    context_aggregator = LLMContextAggregatorPair(context)

    # Create the pipeline
    pipeline = Pipeline([
        transport.input(),
        context_aggregator.user(),
        realtime,
        transport.output(),
        context_aggregator.assistant(),
    ])

    # Create and run the task
    task = PipelineTask(
        pipeline,
        params=PipelineParams(
            allow_interruptions=True,
            enable_metrics=True,
            enable_usage_metrics=True,
        )
    )
    
    # Queue an initial LLM run to start the presentation when pipeline is ready
    @task.event_handler("on_first_participant_joined")
    async def on_first_participant_joined(task, participant):
        logger.info("🎙️ Audio started - beginning presentation...")
        from pipecat.frames.frames import LLMFullResponseStartFrame
        await task.queue_frames([LLMFullResponseStartFrame()])

    logger.info("✅ Voice agent ready!")
    logger.info("📊 PowerPoint viewer window should be open")
    logger.info("🎤 Listening on your microphone...")
    logger.info("💡 The AI will present each slide and ask for questions")
    logger.info("💡 Say 'next slide' to move forward")
    logger.info("💡 Say 'previous slide' to go back")
    logger.info("💡 Ask questions about the content anytime")
    logger.info("🛑 Press Ctrl+C to stop")
    print()

    runner = PipelineRunner()
    await runner.run(task)


if __name__ == "__main__":
    try:
        asyncio.run(main())
    except KeyboardInterrupt:
        logger.info("👋 Shutting down voice agent...")
    except Exception as e:
        logger.error(f"❌ Error: {e}")