"""Tkinter UI for displaying PowerPoint slides."""

import threading
import tkinter as tk
from tkinter import ttk, scrolledtext
from pathlib import Path
from typing import List, Dict, Any, Optional

from pptx import Presentation
from loguru import logger

from models import SlideData


class PowerPointViewerUI:
    """Tkinter UI for displaying PowerPoint slides with text content."""
    
    def __init__(self, pptx_file: str = "./slides.pptx"):
        self.pptx_file = Path(pptx_file)
        self.current_index = 0
        self.slides: List[SlideData] = []
        self.root: Optional[tk.Tk] = None
        self.title_label: Optional[tk.Label] = None
        self.content_text: Optional[scrolledtext.ScrolledText] = None
        self.info_label: Optional[tk.Label] = None
        
        # Load slides
        self._load_slides()
    
    def _extract_text_from_shape(self, shape) -> str:
        """Extract text from a shape."""
        if hasattr(shape, "text"):
            return shape.text.strip()
        elif hasattr(shape, "text_frame"):
            return "\n".join([
                paragraph.text.strip() 
                for paragraph in shape.text_frame.paragraphs
                if paragraph.text.strip()
            ])
        return ""
    
    def _load_slides(self) -> None:
        """Load all slides from the PowerPoint file."""
        if not self.pptx_file.exists():
            logger.error(f"PowerPoint file not found: {self.pptx_file}")
            return
        
        try:
            prs = Presentation(str(self.pptx_file))
            logger.info(f"Loading presentation with {len(prs.slides)} slides")
            
            for idx, slide in enumerate(prs.slides, 1):
                slide_data = self._parse_slide(idx, slide)
                self.slides.append(slide_data)
                logger.info(f"Loaded slide {idx}: {slide_data.title}")
            
            logger.info(f"Successfully loaded {len(self.slides)} slides")
            
        except Exception as e:
            logger.error(f"Error loading PowerPoint: {e}")
    
    def _parse_slide(self, idx: int, slide) -> SlideData:
        """Parse a single slide and return SlideData."""
        # Extract title
        title = "Untitled Slide"
        if slide.shapes.title:
            title = slide.shapes.title.text.strip() or f"Slide {idx}"
        
        # Extract all text content
        content_parts = []
        raw_text_parts = []
        
        for shape in slide.shapes:
            text = self._extract_text_from_shape(shape)
            if text:
                raw_text_parts.append(text)
                # Skip title in content (avoid duplication)
                if shape != slide.shapes.title:
                    content_parts.append(text)
        
        content = "\n\n".join(content_parts)
        raw_text = "\n\n".join(raw_text_parts)
        
        return SlideData(
            slide_number=idx,
            title=title,
            content=content if content else "No content",
            raw_text=raw_text if raw_text else "No text found"
        )
    
    def start(self) -> None:
        """Start the Tkinter UI in a separate thread."""
        thread = threading.Thread(target=self._run_ui, daemon=True)
        thread.start()
    
    def _run_ui(self) -> None:
        """Run the Tkinter main loop."""
        self.root = tk.Tk()
        self._configure_window()
        self._create_widgets()
        
        # Display first slide if available
        if self.slides:
            self.display_current_slide()
        
        self.root.mainloop()
    
    def _configure_window(self) -> None:
        """Configure the main window."""
        self.root.title("Voice Controlled PowerPoint Viewer")
        self.root.geometry("1000x700")
        self.root.configure(bg="#2b2b2b")
    
    def _create_widgets(self) -> None:
        """Create all UI widgets."""
        self._create_app_title()
        self._create_slide_title_frame()
        self._create_content_frame()
        self._create_info_label()
        self._create_navigation_buttons()
    
    def _create_app_title(self) -> None:
        """Create the application title."""
        app_title = tk.Label(
            self.root,
            text="🎤 Voice Controlled PowerPoint Viewer",
            font=("Arial", 16, "bold"),
            bg="#2b2b2b",
            fg="white"
        )
        app_title.pack(pady=10)
    
    def _create_slide_title_frame(self) -> None:
        """Create the slide title frame."""
        title_frame = tk.Frame(self.root, bg="#1e1e1e", relief=tk.RAISED, borderwidth=2)
        title_frame.pack(padx=20, pady=(5, 10), fill=tk.X)
        
        self.title_label = tk.Label(
            title_frame,
            text="No slides loaded",
            font=("Arial", 20, "bold"),
            bg="#1e1e1e",
            fg="#4ECDC4",
            wraplength=900,
            justify=tk.LEFT,
            padx=20,
            pady=15
        )
        self.title_label.pack(fill=tk.X)
    
    def _create_content_frame(self) -> None:
        """Create the content display frame."""
        content_frame = tk.Frame(self.root, bg="#1e1e1e", relief=tk.SUNKEN, borderwidth=2)
        content_frame.pack(padx=20, pady=5, fill=tk.BOTH, expand=True)
        
        self.content_text = scrolledtext.ScrolledText(
            content_frame,
            wrap=tk.WORD,
            font=("Arial", 12),
            bg="#1e1e1e",
            fg="white",
            padx=20,
            pady=20,
            relief=tk.FLAT,
            state=tk.DISABLED
        )
        self.content_text.pack(fill=tk.BOTH, expand=True)
    
    def _create_info_label(self) -> None:
        """Create the slide info label."""
        self.info_label = tk.Label(
            self.root,
            text="No slides loaded",
            font=("Arial", 10),
            bg="#2b2b2b",
            fg="#888888"
        )
        self.info_label.pack(pady=5)
    
    def _create_navigation_buttons(self) -> None:
        """Create navigation buttons."""
        button_frame = tk.Frame(self.root, bg="#2b2b2b")
        button_frame.pack(pady=10)
        
        prev_btn = ttk.Button(button_frame, text="← Previous", command=self.previous_slide)
        prev_btn.pack(side=tk.LEFT, padx=5)
        
        next_btn = ttk.Button(button_frame, text="Next →", command=self.next_slide)
        next_btn.pack(side=tk.LEFT, padx=5)
    
    def display_current_slide(self) -> None:
        """Display the current slide."""
        if not self.slides:
            return
            
        try:
            slide = self.slides[self.current_index]
            
            if self.root:
                self.title_label.config(text=slide.title)
                
                # Update content
                self.content_text.config(state=tk.NORMAL)
                self.content_text.delete(1.0, tk.END)
                self.content_text.insert(1.0, slide.content)
                self.content_text.config(state=tk.DISABLED)
                
                # Update info
                self.info_label.config(
                    text=f"Slide {slide.slide_number} of {len(self.slides)}"
                )
                
            logger.info(f"Displaying slide {slide.slide_number}: {slide.title}")
            
        except Exception as e:
            logger.error(f"Error displaying slide: {e}")
    
    def next_slide(self) -> Optional[Dict[str, Any]]:
        """Move to the next slide."""
        if not self.slides:
            return None
            
        self.current_index = (self.current_index + 1) % len(self.slides)
        self.display_current_slide()
        return self.get_current_slide_info()
    
    def previous_slide(self) -> Optional[Dict[str, Any]]:
        """Move to the previous slide."""
        if not self.slides:
            return None
            
        self.current_index = (self.current_index - 1) % len(self.slides)
        self.display_current_slide()
        return self.get_current_slide_info()
    
    def goto_slide(self, slide_number: int) -> Dict[str, Any]:
        """Go to a specific slide by number (1-based)."""
        if not self.slides:
            return {"error": "No slides available"}
            
        # Convert to 0-based index
        index = slide_number - 1
        
        if 0 <= index < len(self.slides):
            self.current_index = index
            self.display_current_slide()
            return self.get_current_slide_info()
        else:
            return {"error": f"Invalid slide number. Please choose between 1 and {len(self.slides)}"}
    
    def get_current_slide_info(self) -> Dict[str, Any]:
        """Get information about the current slide including full text."""
        if not self.slides:
            return {"error": "No slides available"}
            
        slide = self.slides[self.current_index]
        return {
            "slide_number": slide.slide_number,
            "total_slides": len(self.slides),
            "title": slide.title,
            "content": slide.raw_text,  # Full text for AI to read
            "has_next": self.current_index < len(self.slides) - 1,
            "has_previous": self.current_index > 0
        }
